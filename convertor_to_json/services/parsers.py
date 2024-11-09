from abc import ABC, abstractmethod
import logging
import PyPDF2
import docx
import pptx
from .resume_parser import ResumeParser

logger = logging.getLogger(__name__)

class BaseParser(ABC):
    def __init__(self):
        self.resume_parser = ResumeParser()
    
    def parse(self, file_path):
        """Базовый метод для парсинга файла"""
        try:
            text = self._extract_text(file_path)
            return self.resume_parser.parse(text)
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {str(e)}")
            return self._get_empty_result()
    
    @abstractmethod
    def _extract_text(self, file_path):
        """Метод для извлечения текста из файла"""
        pass
        
    def _get_empty_result(self):
        return {
            'raw_text': '',
            'name': '',
            'email': '',
            'phone': '',
            'experience_years': 0,
            'specialization': '',
            'tech_stack': [],
            'hard_skills': [],
            'soft_skills': [],
            'education': {},
            'work_experience': [],
            'languages': ['Русский'],
            'desired_salary': 0,
            'location': {'country': ''},
            'certifications': []
        }

class PDFParser(BaseParser):
    def _extract_text(self, file_path):
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n'
            return text

class DOCParser(BaseParser):
    def _extract_text(self, file_path):
        # Для .doc файлов может потребоваться дополнительная библиотека
        return ''

class DOCXParser(BaseParser):
    def _extract_text(self, file_path):
        doc = docx.Document(file_path)
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])

class PPTXParser(BaseParser):
    def _extract_text(self, file_path):
        prs = pptx.Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

class TXTParser(BaseParser):
    def _extract_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()